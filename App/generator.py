from typing import List
from pptx import Presentation
from App.models import AnalysisResult, SlideContent, SlideType
import os

class PPTGenerator:
    """PPT 生成器"""
    def __init__(self, template_dir: str = "Templates"):
        self.template_dir = template_dir

    def generate(self, result: AnalysisResult, template: str = "blank") -> str:
        """生成 PPT 文件并返回保存路径"""
        
        #读取模板
        prs = Presentation(f"{self.template_dir}/{template}.pptx")

        #按照类型渲染幻灯片
        for content in result.slides:
            slide = self._render_slide(prs, content)
            
        os.makedirs("output", exist_ok=True)
        prs.save(f"output/{result.title}.pptx")
        return f"output/{result.title}.pptx"
    
    def _render_slide(self, prs: Presentation, content: SlideContent):
        """根据内容类型渲染幻灯片"""
        if content.type == SlideType.TITLE:
            return self._render_title_slide(prs, content)
        elif content.type == SlideType.SECTION:
            return self._render_section_slide(prs, content)
        elif content.type == SlideType.PARAGRAPH:
            return self._render_paragraph_slide(prs, content)
        elif content.type == SlideType.LIST:
            return self._render_list_slide(prs, content)
        elif content.type == SlideType.BULLETS:
            return self._render_bullets_slide(prs, content)
    
    def _render_title_slide(self, prs: Presentation, content: SlideContent):
        """渲染总标题页"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = content.title
        # 可选：添加副标题
        if content.bullet_points:
            subtitle = content.bullet_points[0]
            if hasattr(slide.shapes, 'subtitle'):
                slide.shapes.subtitle.text = subtitle
    
    def _render_section_slide(self, prs: Presentation, content: SlideContent):
        """渲染小标题页（章节页）"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # 显示序号 + 标题
        slide.shapes.title.text = content.title
        # 可选：添加简短描述
        if content.bullet_points:
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.add_paragraph().text = content.bullet_points[0]
    
    def _render_paragraph_slide(self, prs: Presentation, content: SlideContent):
        """渲染大段文字页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = content.title
        text_frame = slide.shapes.placeholders[1].text_frame
        for paragraph in content.bullet_points:
            text_frame.add_paragraph().text = paragraph
    
    def _render_list_slide(self, prs: Presentation, content: SlideContent):
        """渲染平行列举页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = content.title
        text_frame = slide.shapes.placeholders[1].text_frame
        for item in content.bullet_points:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
    
    def _render_bullets_slide(self, prs: Presentation, content: SlideContent):
        """渲染要点列表页"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = content.title
        text_frame = slide.shapes.placeholders[1].text_frame
        # 使用序号（如果有的话）
        for i, point in enumerate(content.bullet_points):
            p = text_frame.add_paragraph()
            # 如果有对应的序号，使用序号；否则直接使用内容
            if i < len(content.number):
                p.text = f"{content.number[i]}. {point}"
            else:
                p.text = f"• {point}"



    def list_templates(self) -> List[str]:
        """列出所有可用模板"""
        pass




